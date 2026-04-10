"""
pptx_service.py - Public entry point for PowerPoint generation.

Orchestrates slide creation and hyperlink wiring.
Sub-modules:
  pptx_constants.py    — colours, geometry, typography
  pptx_helpers.py      — XML/cell helpers, merge utilities
  pptx_data_builder.py — data preparation and pagination
  pptx_slide_builder.py — slide creation functions
"""

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from services.pptx_constants import SLIDE_W, SLIDE_H, NS
from services.pptx_data_builder import build_all_blocks, paginate_blocks
from services.pptx_slide_builder import (
    add_title_slide, add_overview_slide, add_content_slide
)


# ── Public API ────────────────────────────────────────────────────────────────

def generate_summary_pptx(output_path, summary_data, printer, variant,
                           sub_assembly, year, quarter, overview=None):
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    if overview:
        add_title_slide(prs, overview, printer, variant, sub_assembly, year, quarter)

    if overview:
        add_overview_slide(prs, overview, summary_data,
                           printer, variant, sub_assembly, year, quarter)

    all_blocks = build_all_blocks(summary_data)
    slides     = paginate_blocks(all_blocks)

    for slide_blocks in slides:
        add_content_slide(prs, slide_blocks)

    _apply_overview_hyperlinks(prs)

    prs.save(output_path)
    print(f"✓ Summary PowerPoint saved: {output_path}")


# ── Hyperlink post-processing ─────────────────────────────────────────────────

def _apply_overview_hyperlinks(prs):
    """Wire up hyperlinks in the overview remarks cells to content slides."""
    if not hasattr(prs, "_overview_hyperlinks"):
        return

    # Map heading text → slide index
    slide_map = {}
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt and "Input Tray:" in txt and "Category:" in txt:
                    slide_map[txt] = idx

    for link_info in prs._overview_hyperlinks:
        cat_name = link_info["cat_name"]

        target_idx = None
        for heading, idx in slide_map.items():
            if f"Category: {cat_name}" in heading:
                target_idx = idx
                break

        if target_idx is None:
            continue

        cell   = link_info["table"].cell(link_info["row"], link_info["col"])
        tc     = cell._tc
        tf_el  = tc.find(f"{{{NS}}}txBody")
        if tf_el is None:
            continue

        target_slide = prs.slides[target_idx]
        rId = link_info["slide_obj"].part.relate_to(
            target_slide.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
        )

        NS_A = NS
        for para in tf_el.findall(f"{{{NS_A}}}p"):
            runs = para.findall(f"{{{NS_A}}}r")
            if not runs:
                continue
            for run_el in runs:
                rPr = run_el.find(f"{{{NS_A}}}rPr")
                if rPr is None:
                    rPr = etree.SubElement(run_el, f"{{{NS_A}}}rPr")
                    run_el.insert(0, rPr)
                for hl in rPr.findall(f"{{{NS_A}}}hlinkClick"):
                    rPr.remove(hl)
                etree.SubElement(
                    rPr,
                    f"{{{NS_A}}}hlinkClick",
                    attrib={
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id": rId,
                        "action": "ppaction://hlinksldjump"
                    }
                )
                rPr.set("u", "sng")
                clr_el = etree.SubElement(rPr, f"{{{NS_A}}}solidFill")
                etree.SubElement(clr_el, f"{{{NS_A}}}srgbClr",
                                 attrib={"val": "1F4E79"})
