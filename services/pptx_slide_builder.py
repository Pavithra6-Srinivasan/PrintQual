"""
pptx_slide_builder.py - Slide creation functions.

Handles title slide, overview slide, and content slides.
"""

import math
from datetime import datetime

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from services.pptx_constants import (
    MARGIN, TABLE_X, TABLE_W, SLIDE_TOP, SLIDE_BOTTOM,
    FONT, PT_HDG, PT_HEADER, PT_BODY, PT_CONTENT_BODY,
    COL_WHITE, COL_BLACK, COL_BLUE, COL_NAVY,
    COL_PASS_BG, COL_PASS_FG, COL_FAIL_BG, COL_FAIL_FG,
    COL_W, HEADERS, N_COLS, HDR_H, BLOCK_HDG_H, GAP_HDG_TABLE, BLOCK_GAP,
)
from services.pptx_helpers import (
    clear_table_style, set_border, set_text, apply_vertical_merges
)
from services.pptx_data_builder import build_category_results, estimate_row_h

from pptx.dml.color import RGBColor


# ── Title slide ───────────────────────────────────────────────────────────────

def add_title_slide(prs, overview, printer, variant, sub_assembly, year, quarter):
    """
    Slide 1: Clean title slide.
    Line 1: Product + Variant (variant removed from printer if present) — large bold, centred
    Line 2: "{test_condition}   {sub_assembly}   Q{quarter}FY{year}" — navy, centred
    Line 3: test name — black, centred
    Line 4: date — navy, centred
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COL_WHITE

    def add_text(text, y, h, size, color=COL_BLACK, bold=False):
        txb = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(y), Inches(TABLE_W), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
        run.font.name      = FONT

    test_name      = overview.get("objective", "")
    test_condition = overview.get("test_condition", "")

    # Remove variant suffix from printer if it's already there (case-insensitive)
    product_base = printer or ""
    if variant and product_base.lower().endswith(variant.lower()):
        product_base = product_base[:-len(variant)].strip()

    # Title: product base + variant, with proper title case — skip empty parts
    title_line = " ".join(w.capitalize() for w in " ".join(
        p for p in [product_base, variant] if p
    ).split())

    project_phase = overview.get("project_phase", "")
    phase_or_quarter = project_phase if project_phase else f"Q{quarter}FY{year}"

    # Always show sub_assembly so ADF Scan / ADF Copy is unambiguous when the
    # data was split. For non-split cases (e.g. plain "ADF") it is still shown
    # so the slide is self-contained.
    line3 = "   ".join(p for p in [test_condition, sub_assembly, test_name] if p)

    add_text(title_line,
             0.50, 0.65, 30, bold=True)
    add_text(phase_or_quarter,
             1.25, 0.40, 16, color=COL_NAVY)
    add_text(line3,
             1.75, 0.55, 14)
    add_text(f"Date : {datetime.now().strftime('%d %b %Y')}",
             2.40, 0.35, 12, color=COL_NAVY)


# ── Overview slide ────────────────────────────────────────────────────────────

def add_overview_slide(prs, overview, summary_data, printer, variant,
                       sub_assembly, year, quarter):
    """
    Slide 2: Overview with test info table and Key Takeaways table.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COL_WHITE

    test_name      = overview.get("objective", "")
    test_condition = overview.get("test_condition", "")

    # Remove variant suffix from printer if it's already there (case-insensitive)
    product_base = printer or ""
    if variant and product_base.lower().endswith(variant.lower()):
        product_base = product_base[:-len(variant)].strip()
    product_variant = " ".join(p for p in [product_base, variant] if p)

    project_phase = overview.get("project_phase", "")
    phase_or_quarter = project_phase if project_phase else f"Q{quarter}FY{year}"

    # Skip product_variant if it already appears in one of the other parts
    other_parts = [phase_or_quarter, test_condition, sub_assembly, test_name]
    other_text  = " ".join(p for p in other_parts if p).lower()
    pv_lower    = product_variant.lower()
    show_product = product_variant and pv_lower not in other_text

    # Skip sub_assembly if it already appears in any other part
    sa_lower_s2   = (sub_assembly or "").lower()
    other_no_sa   = " ".join(p for p in [
        product_variant if show_product else "",
        phase_or_quarter, test_condition, test_name
    ] if p).lower()
    show_sa_s2 = bool(sub_assembly) and sa_lower_s2 not in other_no_sa

    title_str = " ".join(
        p for p in [
            product_variant if show_product else "",
            phase_or_quarter, test_condition,
            sub_assembly if show_sa_s2 else "",
            test_name, "Summary"
        ]
        if p
    )

    # Title text box
    txb = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(0.05), Inches(TABLE_W), Inches(0.30))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text           = title_str
    run.font.size      = Pt(PT_HDG)
    run.font.bold      = True
    run.font.color.rgb = COL_BLACK
    run.font.name      = FONT

    # Overview info table
    info_table_bottom = _add_overview_info_table(slide, overview)

    # Key Takeaways
    _add_key_takeaways(prs, slide, summary_data, info_table_bottom)


def _add_overview_info_table(slide, overview):
    units      = overview.get("unit_names", [])
    unit_str   = ", ".join(units)
    unit_count = overview.get("unit_count", len(units))
    actual_life = overview.get("actual_life", "")

    headers = [
        "Total Sample Size",
        "Unit Number",
        "Actual Test Life Per Unit",
        "Test Start Date",
        "Test End Date",
    ]
    values = [
        str(unit_count),
        unit_str,
        f"{int(actual_life):,}" if actual_life else "",
        overview.get("test_start", ""),
        overview.get("test_end", ""),
    ]
    col_widths = [1.1, 3.6, 1.8, 1.5, 1.5]   # sums to 9.5"

    # Fixed header height; value height estimated dynamically from unit number wrapping
    hdr_row_h = 0.30   # "Total Sample Size" wraps to 2 lines in narrow column
    tbl_y     = 0.45

    # Estimate lines needed for unit number cell (widest-varying column)
    unit_col_w      = 3.6          # inches (col_widths[1])
    char_w_9pt      = 0.052        # approx inch-per-char for 9pt Calibri
    chars_per_line  = max(1, int(unit_col_w / char_w_9pt))
    unit_lines      = max(1, math.ceil(len(unit_str) / chars_per_line)) if unit_str else 1
    val_row_h       = max(0.35, unit_lines * 0.14)

    tbl_h     = hdr_row_h + val_row_h

    table = slide.shapes.add_table(
        2, len(headers),
        Inches(TABLE_X), Inches(tbl_y),
        Inches(TABLE_W), Inches(tbl_h)
    ).table

    clear_table_style(table)

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    # Set explicit row heights so PowerPoint cannot auto-expand them
    table.rows[0].height = Inches(hdr_row_h)
    table.rows[1].height = Inches(val_row_h)

    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_BLUE
        set_text(cell, hdr, bold=True, fg=COL_WHITE,
                 size=PT_HEADER, align=PP_ALIGN.CENTER)
        set_border(cell)

    for ci, val in enumerate(values):
        cell = table.cell(1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_WHITE
        set_text(cell, val, size=PT_BODY, align=PP_ALIGN.CENTER)
        set_border(cell)

    return tbl_y + tbl_h  # actual bottom y for caller to use


def _add_key_takeaways(prs, slide, summary_data, info_table_bottom):
    """
    Add Key Takeaways tables. Try to fit all in 1 slide (2x2 grid for 4 tables).
    If more than 4 or too tall, use 2 per slide and create additional slides.
    """
    result_groups = build_category_results(summary_data)

    if not hasattr(prs, "_overview_hyperlinks"):
        prs._overview_hyperlinks = []

    n_groups = len(result_groups)

    # Try to fit all groups in one slide (up to 4 in a 2x2 grid)
    if n_groups <= 4:
        label_h = 0.14
        row_gap = 0.35
        # Height of each grid row = max of the two tables' heights
        grid_row_hs = {}
        for idx, g in enumerate(result_groups):
            ri = idx // 2
            h  = _kt_table_h(g["rows"])
            grid_row_hs[ri] = max(grid_row_hs.get(ri, 0), h)
        total_h_needed = 0.22 + sum(
            label_h + KT_LABEL_GAP + h + row_gap
            for h in grid_row_hs.values()
        )
        avail_h = SLIDE_BOTTOM - (info_table_bottom + 0.12)

        if total_h_needed <= avail_h:
            _render_kt_grid(slide, prs, result_groups, info_table_bottom + 0.12, TABLE_W)
            return

    # Fallback: 2 tables per slide
    GROUPS_PER_SLIDE = 2

    for slide_idx, start_idx in enumerate(range(0, len(result_groups), GROUPS_PER_SLIDE)):
        groups_for_slide = result_groups[start_idx:start_idx + GROUPS_PER_SLIDE]

        # Use the passed slide for first batch, create new slides for overflow
        current_slide = slide if slide_idx == 0 else prs.slides.add_slide(prs.slide_layouts[6])
        current_slide.background.fill.solid()
        current_slide.background.fill.fore_color.rgb = COL_WHITE

        # "Key Takeaways" heading (only on first slide)
        if slide_idx == 0:
            title_y = info_table_bottom + 0.12
        else:
            title_y = 0.15

        title_box = current_slide.shapes.add_textbox(
            Inches(TABLE_X), Inches(title_y),
            Inches(TABLE_W), Inches(0.28))
        tf  = title_box.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text           = "Key Takeaways" + (f" (cont.)" if slide_idx > 0 else "")
        run.font.bold      = True
        run.font.underline = True
        run.font.size      = Pt(14)
        run.font.name      = FONT

        if len(groups_for_slide) == 1:
            # Single table — full width
            table_y = title_y + 0.34
            group   = groups_for_slide[0]
            _draw_kt_table(current_slide, prs, group, TABLE_X, table_y, TABLE_W)
        else:
            # Two tables side-by-side
            gap   = 0.10
            tbl_w = (TABLE_W - gap) / 2
            row_height = 0.22
            hdr_height = 0.22
            row_count = max(len(g["rows"]) for g in groups_for_slide)
            tbl_h = hdr_height + row_height * row_count

            start_y = title_y + 0.34

            for idx, group in enumerate(groups_for_slide):
                x = TABLE_X + idx * (tbl_w + gap)

                # Condition label above each table
                lbl = current_slide.shapes.add_textbox(
                    Inches(x), Inches(start_y), Inches(tbl_w), Inches(0.18))
                lp  = lbl.text_frame.paragraphs[0]
                lr  = lp.add_run()
                lr.text      = group["test_condition"] or "Ambient"
                lr.font.bold = True
                lr.font.size = Pt(8)
                lr.font.name = FONT

                _draw_kt_table(current_slide, prs, group, x, start_y + 0.18, tbl_w)


KT_HDR_H      = 0.18   # header row height
KT_PASS_ROW_H = 0.22   # height for rows with no remarks (8pt text + ~0.05" top/bottom cell margins)
KT_REM_LINE_H = 0.10   # additional height per wrapped remark line (8pt Calibri @ 80% lnSpc ≈ 0.089")
KT_LABEL_GAP  = 0.05   # gap between condition label and table
KT_CHAR_W     = 0.048  # approx inch-per-char for 8pt Calibri (used to derive chars_per_line)


def _estimate_remark_lines(remark_text, chars_per_line):
    """Estimate wrapped line count for a remark string."""
    if not remark_text:
        return 1
    return max(1, math.ceil(len(remark_text) / chars_per_line))


def _kt_row_heights(rows, col_w_rem):
    """Estimate per-row heights based on actual remark content and column width."""
    chars_per_line = max(20, int(col_w_rem / KT_CHAR_W))
    heights = []
    for row in rows:
        remarks = row.get("remarks", [])
        if not remarks:
            heights.append(KT_PASS_ROW_H)
        else:
            total_lines = sum(_estimate_remark_lines(r, chars_per_line) for r in remarks)
            heights.append(KT_PASS_ROW_H + KT_REM_LINE_H * max(0, total_lines - 1))
    return heights if heights else [KT_PASS_ROW_H]


def _kt_table_h(rows, col_w_rem=1.35):
    return KT_HDR_H + sum(_kt_row_heights(rows, col_w_rem))


def _render_kt_grid(slide, prs, result_groups, start_y, max_width):
    """Render up to 4 Key Takeaways tables in a 2x2 grid."""
    title_box = slide.shapes.add_textbox(
        Inches(TABLE_X), Inches(start_y),
        Inches(max_width), Inches(0.28))
    tf  = title_box.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text           = "Key Takeaways"
    run.font.bold      = True
    run.font.underline = True
    run.font.size      = Pt(14)
    run.font.name      = FONT

    gap     = 0.08
    tbl_w   = (max_width - gap) / 2
    label_h = 0.14
    row_gap = 0.35

    tbl_y_start = start_y + 0.36

    # Pre-compute table heights per group
    tbl_heights = [_kt_table_h(g["rows"]) for g in result_groups]

    # Height of each grid row = max of the two tables in that row
    grid_row_heights = {}
    for idx in range(len(result_groups)):
        row_idx = idx // 2
        h = tbl_heights[idx]
        grid_row_heights[row_idx] = max(grid_row_heights.get(row_idx, 0), h)

    # Cumulative y offset per grid row
    grid_row_y = {}
    cursor = tbl_y_start
    for row_idx in sorted(grid_row_heights):
        grid_row_y[row_idx] = cursor
        cursor += label_h + KT_LABEL_GAP + grid_row_heights[row_idx] + row_gap

    for idx, group in enumerate(result_groups):
        col_idx = idx % 2
        row_idx = idx // 2

        # If this group is alone in its grid row, expand it to full width
        is_lone = col_idx == 0 and (idx + 1 >= len(result_groups))
        w = max_width if is_lone else tbl_w
        x = TABLE_X + col_idx * (tbl_w + gap)
        y = grid_row_y[row_idx]

        # Condition label
        lbl = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(label_h))
        lp  = lbl.text_frame.paragraphs[0]
        lr  = lp.add_run()
        lr.text      = group["test_condition"] or "Ambient"
        lr.font.bold = True
        lr.font.size = Pt(7)
        lr.font.name = FONT

        _draw_kt_table(slide, prs, group, x, y + label_h + KT_LABEL_GAP, w)


def _draw_kt_table(slide, prs, group, x, y, width):
    """Render a single Key Takeaways table at the given position.

    Columns: Category | Media Type | Media Cat | Result | Remarks
    Category cell is shown only on the first row for each category
    (visual deduplication without actual PowerPoint cell merging).
    """
    rows        = group.get("rows", [])
    col_w_cat   = 0.90
    col_w_res   = 1.30
    col_w_rem   = max(0.5, width - col_w_cat - col_w_res)
    col_widths  = [col_w_cat, col_w_res, col_w_rem]

    row_heights = _kt_row_heights(rows, col_w_rem)
    table_h     = KT_HDR_H + sum(row_heights)

    table = slide.shapes.add_table(
        1 + len(rows), 3,
        Inches(x), Inches(y),
        Inches(width), Inches(table_h)
    ).table

    clear_table_style(table)
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    # Set row heights explicitly
    table.rows[0].height = Inches(KT_HDR_H)
    for ri, h in enumerate(row_heights):
        table.rows[ri + 1].height = Inches(h)

    for ci, hdr in enumerate(["Category", "Result", "Remarks"]):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_BLUE
        set_text(cell, hdr, bold=True, fg=COL_WHITE,
                 size=8, align=PP_ALIGN.CENTER)
        set_border(cell)

    last_cat = None
    for ri, row_info in enumerate(rows):
        tr         = ri + 1
        cat        = row_info["category"]
        media_type = row_info.get("media_type", "")
        media_cat  = row_info.get("media_cat",  "")
        result     = row_info["result"]
        remarks    = row_info["remarks"]

        # Show category name only on the first row for that category
        display_cat = cat if cat != last_cat else ""
        last_cat    = cat

        cell = table.cell(tr, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_WHITE
        set_text(cell, display_cat, size=8)
        set_border(cell)

        result_style = row_info.get("result_style", result)
        cell = table.cell(tr, 1)
        cell.fill.solid()
        if result_style in ("PASS", "No issue"):
            cell.fill.fore_color.rgb = COL_PASS_BG
            set_text(cell, result, bold=True, fg=COL_PASS_FG,
                     size=7, align=PP_ALIGN.LEFT)
        elif result_style in ("Observation", "With observation", "With Observation"):
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x9C)  # amber
            set_text(cell, result, bold=True, fg=RGBColor(0x7F, 0x4E, 0x00),
                     size=8, align=PP_ALIGN.CENTER)
        else:
            cell.fill.fore_color.rgb = COL_FAIL_BG
            set_text(cell, result, bold=True, fg=COL_FAIL_FG,
                     size=8, align=PP_ALIGN.CENTER)
        set_border(cell)

        cell = table.cell(tr, 2)
        set_text(cell, "\n".join(remarks), size=8)
        set_border(cell)

        if remarks:
            prs._overview_hyperlinks.append({
                "slide_obj": slide,
                "table":     table,
                "row":       tr,
                "col":       2,   # remarks is column 2
                "cat_name":  cat,
            })


# ── Content slide ─────────────────────────────────────────────────────────────

def add_content_slide(prs, slide_blocks):
    """Render one slide containing one or more heading+table blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COL_WHITE

    cursor_y = SLIDE_TOP

    for block_idx, block in enumerate(slide_blocks):
        if block_idx > 0:
            cursor_y += BLOCK_GAP

        heading   = block["heading"]
        flat_rows = block["rows"]

        # Heading text box
        txb = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(cursor_y),
            Inches(TABLE_W), Inches(BLOCK_HDG_H))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text           = heading
        run.font.size      = Pt(PT_HDG)
        run.font.bold      = True
        run.font.color.rgb = COL_BLACK
        run.font.name      = FONT

        cursor_y += BLOCK_HDG_H + GAP_HDG_TABLE

        if not flat_rows:
            continue

        # ── Determine active (non-empty) columns ─────────────────────────────
        # Cols 4-7 (Overall Result, Error Type, Media Name, Unit) are always
        # included. Cols 0-3 (Media Type, Media Cat, Print Mode, Spec) are
        # only included if at least one row has a non-empty value.
        ALWAYS_INCLUDE = {4, 5, 6, 7}
        active_cols = [
            ci for ci in range(N_COLS)
            if ci in ALWAYS_INCLUDE or any(row["cols"][ci] for row in flat_rows)
        ]
        orig_to_new  = {orig: new for new, orig in enumerate(active_cols)}
        n_active     = len(active_cols)
        headers_used = [HEADERS[ci] for ci in active_cols]
        widths_used  = [COL_W[ci] for ci in active_cols]

        # Redistribute freed width to the widest remaining column
        freed = TABLE_W - sum(widths_used)
        if freed > 0:
            widest_new = widths_used.index(max(widths_used))
            widths_used[widest_new] += freed

        row_heights = [estimate_row_h(r) for r in flat_rows]
        table_h     = HDR_H + sum(row_heights)
        n_total     = 1 + len(flat_rows)

        tbl_shape = slide.shapes.add_table(
            n_total, n_active,
            Inches(TABLE_X), Inches(cursor_y),
            Inches(TABLE_W), Inches(table_h))
        tbl = tbl_shape.table

        clear_table_style(tbl)

        for new_ci, cw in enumerate(widths_used):
            tbl.columns[new_ci].width = Inches(cw)

        # Header row
        for new_ci, hdr in enumerate(headers_used):
            cell = tbl.cell(0, new_ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COL_BLUE
            set_text(cell, hdr, bold=True, fg=COL_WHITE,
                     size=PT_CONTENT_BODY, align=PP_ALIGN.CENTER)
            set_border(cell)

        # Data rows
        result_new_ci = orig_to_new.get(4)
        for ri, row in enumerate(flat_rows):
            tr = ri + 1
            for new_ci, orig_ci in enumerate(active_cols):
                val  = row["cols"][orig_ci]
                cell = tbl.cell(tr, new_ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COL_WHITE

                if new_ci == result_new_ci and val in ("PASS", "FAIL", "NO SPEC"):
                    if val == "PASS":
                        cell.fill.fore_color.rgb = COL_PASS_BG
                        set_text(cell, val, bold=True, fg=COL_PASS_FG,
                                 size=PT_CONTENT_BODY, align=PP_ALIGN.CENTER)
                    elif val == "FAIL":
                        cell.fill.fore_color.rgb = COL_FAIL_BG
                        set_text(cell, val, bold=True, fg=COL_FAIL_FG,
                                 size=PT_CONTENT_BODY, align=PP_ALIGN.CENTER)
                    else:
                        cell.fill.fore_color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
                        set_text(cell, val, bold=True,
                                 fg=RGBColor(0x40, 0x40, 0x40),
                                 size=PT_CONTENT_BODY, align=PP_ALIGN.CENTER)
                    set_border(cell)
                    continue

                set_text(cell, str(val) if val else "",
                         size=PT_CONTENT_BODY, align=PP_ALIGN.LEFT)
                set_border(cell)

        apply_vertical_merges(tbl, flat_rows, orig_to_new)

        # Set row heights AFTER merges so PowerPoint respects them
        tbl.rows[0].height = Inches(HDR_H)
        for ri, h in enumerate(row_heights):
            tbl.rows[ri + 1].height = Inches(h)

        cursor_y += table_h
